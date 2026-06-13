#!/usr/bin/env python3
"""Spark Planners (in association with TentnMore) — company profile & services deck.

Reuses the themed engine in build_spark_deck.py. Content drawn from TentnMore's
service notes and company profile; gallery photos extracted from that profile
live in public/generated/tentnmore/.

Run:  python3 build_spark_tentnmore_deck.py            # all four themes
      python3 build_spark_tentnmore_deck.py emerald    # one theme
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import build_spark_deck as E  # themed engine (helpers + THEMES)

BASE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "spark-tentnmore-profile")
T = lambda f: os.path.join(E.PUB, "generated", "tentnmore", f)  # TentnMore asset
FOOT = "Spark Planners · in association with TentnMore"

# ---- Service content (condensed from the notes) -----------------------------
PLANNING = [
    ("Destination Wedding Scouting",
     "Breathtaking backdrops countrywide, with guest travel, stay and local legal compliance handled by our concierge."),
    ("Wedding Venue Selection",
     "Ideal spaces matched to you, virtual or in-person site visits, and negotiated premium rates with optimised layouts."),
    ("Event Timeline",
     "A master blueprint with minute-by-minute schedules and built-in buffers to keep every vendor team in sync."),
    ("Logistics, Hospitality & Transport",
     "A premium fleet, airport welcome desks, luggage-tagged check-ins and 24/7 on-site guest support."),
]
CEREMONIES = [
    ("Proposal Planning",
     "Bespoke surprise concepts with secret locations, hidden cameras, music, flowers and flawless timing."),
    ("Mehandi · Haldi · Sangeet",
     "Sacred family rituals with festive, high-octane energy, bold palettes and interactive stalls like lac bangles and chaat."),
    ("Wedding & Reception",
     "Serene ceremony altars, grand couple entries, and a seamless flow from speeches to after-party."),
    ("Choreography",
     "Routines tailored to your skill level, energetic family mashups and stage blocking under the lights."),
]
DESIGN = [
    ("Wedding Décor for All Functions",
     "Custom-fabricated stages and backdrops, immersive smart lighting, and premium linens, drapes and carpets."),
    ("Flowers & Bouquets",
     "Massive floral installations from rare, farm-sourced blooms, plus bouquets, boutonnieres, garlands and scent design."),
    ("Wedding Invitation",
     "Custom stationery suites, interactive digital websites, real-time RSVP tracking and gourmet invite hampers."),
]
FOOD = [
    ("Catering & Wedding Menu",
     "Bespoke multi-cuisine menus shaped through private chef tastings, live theatre food stations, and full vegan, gluten-free and allergy care."),
    ("Wedding Cake",
     "Edible sugar-art masterpieces matched to your dress, premium fillings, dramatic reveals, cupcakes and dessert tables."),
]
EXPERIENCE = [
    ("Photography & Cinematography",
     "Top-tier storytellers for editorial pre-wedding shoots, movie-quality films and fast quick-preview delivery."),
    ("Entertainment",
     "Premium live acts — DJs, brass bands, acoustic duos — plus magicians, flash mobs, photo booths and 360° cameras."),
    ("Artist & Celebrity Management",
     "A-list musicians, actors and hosts with full technical, rider, green-room and private-security handling."),
    ("Favours & Gifts",
     "Luxury takeaway tokens and local souvenirs, custom-monogrammed packaging and curated welcome baskets."),
]


def build_deck(out_path):
    A, TI, TX, MU = E.ACCENT, E.TITLE, E.TEXT, E.MUTED
    CARD, PANEL = E.CARD, E.PANEL
    SW, SH = E.SW, E.SH

    def service_slide(kick, head, items, photo, bg=None):
        s = E.slide(bg)
        E.kicker(s, kick)
        E.title(s, head)
        has_photo = bool(E.prefer(photo))
        lw = 6.25 if has_photo else 11.5
        n = len(items)
        top, gap = 2.15, 0.2
        rowh = (6.45 - top - (n - 1) * gap) / n
        for i, (name, desc) in enumerate(items):
            y = top + i * (rowh + gap)
            E.rect(s, 0.92, y, lw, rowh, CARD, A, 0.6)
            E.rect(s, 0.92, y, 0.07, rowh, A)
            E.txt(s, 1.2, y + 0.13, lw - 0.5, 0.4, name, 14.5, TI, bold=True, font=E.HEAD)
            E.txt(s, 1.2, y + 0.52, lw - 0.5, rowh - 0.5, desc, 11.5, MU, spacing=1.08)
        if has_photo:
            E.pic_cover(s, photo, 7.45, 2.15, 4.95, 4.45)
        E.footer(s, FOOT)
        return s

    # 1 Cover
    s = E.slide()
    E.rect(s, 0.45, 0.45, SW - 0.9, SH - 0.9, None, A, 1.25)
    E.rect(s, 0.62, 0.62, SW - 1.24, SH - 1.24, None, A, 0.5)
    E.txt(s, 0, 1.55, SW, 0.4, "LUXURY WEDDINGS & EVENTS · DESIGN TO DELIVERY", 13.5, A, bold=True, align=PP_ALIGN.CENTER)
    E.rect(s, SW / 2 - 1.1, 2.2, 2.2, 0.025, A)
    E.txt(s, 0, 2.7, SW, 1.2, "SPARK PLANNERS", 56, TI, bold=True, font=E.HEAD, align=PP_ALIGN.CENTER)
    E.txt(s, 0, 3.95, SW, 0.5, "in association with TentnMore Event Planners", 18, TX, font=E.HEAD, italic=True, align=PP_ALIGN.CENTER)
    E.txt(s, 0, 4.75, SW, 0.5, "Designing unforgettable celebrations — plan · design · film · execute.", 15, TX, italic=True, align=PP_ALIGN.CENTER)
    E.rect(s, SW / 2 - 1.1, 5.5, 2.2, 0.025, A)
    E.txt(s, 0, 5.8, SW, 0.4, "COMPANY PROFILE & SERVICES", 12.5, MU, bold=True, align=PP_ALIGN.CENTER)

    # 2 Who we are
    s = E.slide()
    E.kicker(s, "WHO WE ARE")
    E.title(s, "We don't just plan weddings")
    E.txt(s, 0.92, 2.15, 6.5, 3.8,
          "We bring your unique love story to life. Every detail reflects your "
          "personality, style and dreams — so you can manage nothing and enjoy "
          "every single moment.\n\n"
          "Spark Planners pairs design-led, technology-enabled planning with "
          "TentnMore's full-service execution across Delhi/NCR and destinations "
          "nationwide. From first concept to the final frame, one accountable team "
          "owns your celebration.",
          16, TX, spacing=1.18)
    for i, (h, d) in enumerate([("Design-led", "Every element drawn to scale before it is built."),
                                ("Full-service", "Décor, films, catering, logistics — under one roof."),
                                ("On the ground", "Seasoned Delhi/NCR execution, calm and flawless.")]):
        y = 2.2 + i * 1.2
        E.rect(s, 7.75, y, 4.65, 1.05, CARD, A, 0.75)
        E.txt(s, 8.0, y + 0.15, 4.1, 0.4, h, 16, TI, bold=True, font=E.HEAD)
        E.txt(s, 8.0, y + 0.56, 4.1, 0.4, d, 12, MU)
    E.footer(s, FOOT)

    # 3 The partnership
    s = E.slide(PANEL)
    E.kicker(s, "THE PARTNERSHIP")
    E.title(s, "Two studios, one seamless team")
    for i, (h, sub, body) in enumerate([
        ("Spark Planners", "Design & technology",
         "Architecture-grade design of venue, mandap, décor and lighting — drawn to scale, previewed in 3D, and presented as client-ready boards. In-house films capture every moment."),
        ("TentnMore Event Planners", "Execution & hospitality",
         "Delhi/NCR's full-service event house: vendors, fabrication, catering, logistics and on-ground delivery, with a proven portfolio across weddings, parties and corporate events.")]):
        x = 0.92 + i * 6.05
        E.rect(s, x, 2.15, 5.7, 4.1, CARD, A, 0.75)
        E.rect(s, x, 2.15, 0.08, 4.1, A)
        E.txt(s, x + 0.32, 2.4, 5.1, 0.5, h, 21, TI, bold=True, font=E.HEAD)
        E.txt(s, x + 0.32, 2.95, 5.1, 0.35, sub.upper(), 12, A, bold=True)
        E.txt(s, x + 0.32, 3.45, 5.1, 2.6, body, 14, TX, spacing=1.2)
    E.footer(s, FOOT)

    # 4 What we do (overview)
    s = E.slide()
    E.kicker(s, "WHAT WE DO")
    E.title(s, "Full-service, end to end")
    cats = ["Planning & Coordination", "Ceremonies & Functions", "Design, Décor & Florals",
            "Catering, Menu & Cake", "Films & Photography", "Entertainment & Choreography",
            "Guest Experience & Logistics", "Beyond Weddings"]
    cw, ch, gx, gy = 2.78, 1.65, 0.26, 0.3
    for i, c in enumerate(cats):
        r, col = divmod(i, 4)
        x = 0.92 + col * (cw + gx); y = 2.2 + r * (ch + gy)
        E.rect(s, x, y, cw, ch, CARD, A, 0.75)
        E.rect(s, x, y, cw, 0.08, A)
        E.txt(s, x + 0.22, y + 0.3, cw - 0.44, 1.1, c, 15, TI, bold=True, font=E.HEAD, spacing=1.05)
    E.footer(s, FOOT)

    # 5-9 Service detail slides
    service_slide("PLANNING & COORDINATION", "From scouting to seamless logistics", PLANNING, T("acc-venue-wide.png"), PANEL)
    service_slide("CEREMONIES & FUNCTIONS", "Every ritual, beautifully staged", CEREMONIES, T("g7-outdoor-stage.png"))
    service_slide("DESIGN, DÉCOR & FLORALS", "Visual storytelling, drawn to scale", DESIGN, T("g2-floral-wall.png"), PANEL)
    service_slide("CATERING, MENU & CAKE", "A dining experience to remember", FOOD, T("g4-golden-dining.png"))
    service_slide("FILMS, ENTERTAINMENT & GUEST EXPERIENCE", "The day, captured and celebrated", EXPERIENCE, T("acc-reception.png"), PANEL)

    # 10 Beyond weddings
    s = E.slide()
    E.kicker(s, "BEYOND WEDDINGS")
    E.title(s, "Celebrations of every kind")
    for i, (h, d) in enumerate([
        ("Birthday Party Decorations",
         "Sparkling, colourful set-ups that turn any milestone into a stunning, memorable celebration for the whole family."),
        ("Corporate Event Management",
         "Complete corporate planning — styling, theming, entertainment and audio-visual — delivered end to end, on brief and on time.")]):
        x = 0.92 + i * 6.05
        E.rect(s, x, 2.2, 5.7, 2.5, CARD, A, 0.75)
        E.rect(s, x, 2.2, 0.08, 2.5, A)
        E.txt(s, x + 0.32, 2.5, 5.1, 0.5, h, 18, TI, bold=True, font=E.HEAD)
        E.txt(s, x + 0.32, 3.2, 5.1, 1.4, d, 13.5, TX, spacing=1.2)
    E.txt(s, 0.92, 5.15, 11.5, 0.4, "Also: cocktail evenings · bachelor & bachelorette parties · anniversaries · private soirées.", 13.5, MU, italic=True)
    E.footer(s, FOOT)

    # 11 Our edge (platform)
    s = E.slide(PANEL)
    E.kicker(s, "OUR EDGE")
    E.title(s, "Powered by a Wedding Operating System")
    E.txt(s, 0.92, 2.0, 11.5, 0.7,
          "What you approve is exactly what we build. Every Spark celebration is planned on a "
          "connected design platform — not sketched on napkins.", 15.5, TX, spacing=1.1)
    for i, (n, l) in enumerate([("156", "design surfaces"), ("1,441", "knowledge nodes"),
                                ("4,539", "linked relations"), ("17", "board pages")]):
        x = 0.92 + i * 3.0
        E.rect(s, x, 3.05, 2.7, 1.7, CARD, A, 0.75)
        E.txt(s, x, 3.3, 2.7, 0.9, n, 44, TI, bold=True, font=E.HEAD, align=PP_ALIGN.CENTER)
        E.txt(s, x, 4.25, 2.7, 0.4, l, 12.5, MU, align=PP_ALIGN.CENTER)
    for i, t in enumerate(["AI design co-pilot for ideas & options", "3D / CAD studios for true-to-scale design",
                           "Board Composer for client-ready presentations", "Vedi & Muhurat intelligence built in"]):
        E.txt(s, 0.95, 5.15 + i * 0.42, 11.5, 0.4, "•  " + t, 13.5, TX)
    E.footer(s, FOOT)

    # 12 How we work
    s = E.slide()
    E.kicker(s, "HOW WE WORK")
    E.title(s, "From first call to final frame")
    steps = [("01", "Discover", "Vision, guest count, dates, budget."),
             ("02", "Design", "Venue, mandap, décor and lighting to scale."),
             ("03", "Plan", "Vendors, timeline, logistics and muhurat."),
             ("04", "Produce", "Fabrication, crews, call sheets, delivery."),
             ("05", "Film", "Pre-wedding through full wedding film."),
             ("06", "Celebrate", "We run the day; you live it.")]
    for i, (n, h, d) in enumerate(steps):
        r, c = divmod(i, 3)
        x = 0.92 + c * 4.05; y = 2.2 + r * 1.9
        E.rect(s, x, y, 3.8, 1.6, CARD, A, 0.75)
        E.txt(s, x + 0.25, y + 0.16, 1.5, 0.6, n, 30, A, bold=True, font=E.HEAD)
        E.txt(s, x + 1.35, y + 0.22, 2.3, 0.4, h, 16, TI, bold=True, font=E.HEAD)
        E.txt(s, x + 0.25, y + 0.78, 3.3, 0.7, d, 12.5, MU)
    E.footer(s, FOOT)

    # 13 Gallery
    s = E.slide(PANEL)
    E.kicker(s, "SELECTED WORK")
    E.title(s, "Recent celebrations")
    gx0, gy0, gw, gh, ggx, ggy = 0.92, 2.15, 2.72, 2.05, 0.19, 0.22
    photos = ["g1-reception-stage.png", "g2-floral-wall.png", "g3-banquet.png", "g4-golden-dining.png",
              "g5-lawn-arch.png", "g6-mandap-night.png", "g7-outdoor-stage.png", "g8-floral-lights.png"]
    for i, p in enumerate(photos):
        r, c = divmod(i, 4)
        x = gx0 + c * (gw + ggx); y = gy0 + r * (gh + ggy)
        E.pic_cover(s, T(p), x, y, gw, gh)
    E.footer(s, FOOT)

    # 14 Why us
    s = E.slide()
    E.kicker(s, "WHY US")
    E.title(s, "The wow factor, end to end")
    E.txt(s, 1.4, 2.5, 10.5, 2.6,
          "“We aim to provide every detail touch to your rituals and imbibe the wow "
          "factor — to make it look simple yet elegant.”",
          24, TI, italic=True, font=E.HEAD, align=PP_ALIGN.CENTER, spacing=1.2)
    E.txt(s, 1.4, 5.0, 10.5, 0.8,
          "One team. One accountable plan. A flawless, unforgettable celebration "
          "from start to finish.", 15, TX, align=PP_ALIGN.CENTER, spacing=1.15)
    E.footer(s, FOOT)

    # 15 Contact / Thank you
    s = E.slide()
    E.rect(s, 0.45, 0.45, SW - 0.9, SH - 0.9, None, A, 1.25)
    E.rect(s, 0.62, 0.62, SW - 1.24, SH - 1.24, None, A, 0.5)
    E.rect(s, SW / 2 - 1.1, 2.3, 2.2, 0.025, A)
    E.txt(s, 0, 2.75, SW, 1.0, "Let's create your story", 40, TI, bold=True, font=E.HEAD, align=PP_ALIGN.CENTER)
    E.txt(s, 0, 4.0, SW, 0.45, "SPARK PLANNERS  ·  Sanjeeta Khaira", 17, TX, bold=True, align=PP_ALIGN.CENTER)
    E.txt(s, 0, 4.5, SW, 0.4, "in association with TentnMore Event Planners", 13.5, MU, italic=True, align=PP_ALIGN.CENTER)
    E.txt(s, 0, 5.15, SW, 0.4, "tentnmore@gmail.com    ·    +91 95402 30027    ·    www.tentnmore.in", 14, TI, align=PP_ALIGN.CENTER)
    E.rect(s, SW / 2 - 1.1, 5.75, 2.2, 0.025, A)

    E.prs.save(out_path)


def main():
    keys = sys.argv[1:] or list(E.THEMES)
    for key in keys:
        if key not in E.THEMES:
            print(f"unknown theme '{key}'. options: {', '.join(E.THEMES)}"); continue
        t = E.THEMES[key]
        E.apply_theme(t)
        E.prs = Presentation()
        E.prs.slide_width = Inches(E.SW); E.prs.slide_height = Inches(E.SH)
        E.BLANK = E.prs.slide_layouts[6]
        out = BASE + t["suffix"] + ".pptx"
        build_deck(out)
        print(f"saved {os.path.basename(out)}  ({t['label']}, {len(E.prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
