#!/usr/bin/env python3
"""Build the Spark Planners client presentation (.pptx) in multiple themes.

Run:  python3 build_spark_deck.py            # builds all themes
      python3 build_spark_deck.py emerald    # builds one theme by key

Themes share one layout; only the palette changes. Photoreal assets in
public/generated/ (see generate_assets.py) are theme-neutral and used by
every theme when present, with a clean fallback when absent.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def C(hexstr):
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


PUB = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "public")
BASE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "spark-planners-client-deck")

# ----- Themes: semantic tokens ------------------------------------------------
# BG=background, PANEL/CARD=surfaces, ACCENT=gold rules, TITLE=heading text,
# TEXT=body text, MUTED=secondary text, ON_ACCENT=text drawn on an ACCENT fill.
THEMES = {
    "midnight-gold": dict(label="Midnight & Gold", suffix="", BG="02100E", PANEL="0A1F1B",
        CARD="102420", ACCENT="D9AA46", TITLE="F8D78B", TEXT="FFF6DF", MUTED="C8B98E", ON_ACCENT="02100E"),
    "royal-emerald": dict(label="Royal Emerald", suffix="-emerald", BG="04231A", PANEL="06321F",
        CARD="0A3B29", ACCENT="E0B65A", TITLE="F6DC97", TEXT="F2EEDD", MUTED="AEC1A6", ON_ACCENT="04231A"),
    "royal-maroon": dict(label="Royal Maroon", suffix="-maroon", BG="2A0A12", PANEL="390F1A",
        CARD="481524", ACCENT="E5C26A", TITLE="F7E3A6", TEXT="FBF1E6", MUTED="D2B3A4", ON_ACCENT="2A0A12"),
    "ivory-gold": dict(label="Ivory & Gold", suffix="-ivory", BG="FBF5E9", PANEL="F1E6D2",
        CARD="F6EEDC", ACCENT="B6862E", TITLE="6E2433", TEXT="3B2E22", MUTED="8C7B62", ON_ACCENT="FBF5E9"),
}

# Literal floral-swatch colours (true to the flowers; same in every theme).
LOTUS, MARIGOLD, JASMINE, NEEM = C("C75B86"), C("E2952A"), C("FBF7EC"), C("6E8B4A")
HEAD, BODY = "Georgia", "Calibri"
SW, SH = 13.333, 7.5

# Palette globals (set per theme by apply_theme)
BG = PANEL = CARD = ACCENT = TITLE = TEXT = MUTED = ON_ACCENT = C("000000")
prs = None
BLANK = None


def apply_theme(t):
    global BG, PANEL, CARD, ACCENT, TITLE, TEXT, MUTED, ON_ACCENT
    BG, PANEL, CARD = C(t["BG"]), C(t["PANEL"]), C(t["CARD"])
    ACCENT, TITLE, TEXT = C(t["ACCENT"]), C(t["TITLE"]), C(t["TEXT"])
    MUTED, ON_ACCENT = C(t["MUTED"]), C(t["ON_ACCENT"])


# ----- Drawing helpers (resolve colours at call time) -------------------------
def slide(bg=None):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG if bg is None else bg
    return s


def rect(s, l, t, w, h, fill, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, l, t, w, h, text, size, color=None, bold=False, font=None, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, spacing=None):
    color = TEXT if color is None else color
    font = BODY if font is None else font
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = color
    return tb


def pic_cover(s, path, l, t, w, h, focus=0.5):
    """focus biases the cropped axis: 0=keep top/left, 0.5=center, 1=keep bottom/right."""
    if not path or not os.path.exists(path):
        rect(s, l, t, w, h, CARD, ACCENT, 0.75)
        return
    from PIL import Image
    iw, ih = Image.open(path).size
    box_ar = w / h; img_ar = iw / ih
    pic = s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    if img_ar > box_ar:
        crop = 1 - box_ar / img_ar; pic.crop_left = crop * focus; pic.crop_right = crop * (1 - focus)
    else:
        crop = 1 - img_ar / box_ar; pic.crop_top = crop * focus; pic.crop_bottom = crop * (1 - focus)
    pic.line.color.rgb = ACCENT; pic.line.width = Pt(1)
    return pic


def overlay(s, l, t, w, h, alpha=60):
    sp = rect(s, l, t, w, h, BG)
    srgb = sp.fill.fore_color._xFill.find(qn("a:srgbClr"))
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))}))
    return sp


def kicker(s, text, l=0.92, t=0.62):
    txt(s, l, t, 8, 0.3, text, 12, ACCENT, bold=True)


def title(s, text, l=0.92, t=0.95, w=11.5, size=40):
    txt(s, l, t, w, 1.1, text, size, TITLE, bold=True, font=HEAD)


def footer(s, text="Spark Planners · Sanjeeta Khaira"):
    txt(s, 0.92, 7.02, 11.5, 0.3, text, 9, MUTED)


P = lambda f: os.path.join(PUB, f)
G = lambda f: os.path.join(PUB, "generated", f)


def prefer(*paths):
    for p in paths:
        if p and os.path.exists(p):
            return p
    return None


# ----- Deck construction ------------------------------------------------------
def build_deck(out_path):
    # 1 Cover
    s = slide()
    _hero = prefer(G("udaipur-palace-night-hero.png"))
    if _hero:
        pic_cover(s, _hero, 0, 0, SW, SH); overlay(s, 0, 0, SW, SH, 64)
    rect(s, 0.45, 0.45, SW - 0.9, SH - 0.9, None, ACCENT, 1.25)
    rect(s, 0.62, 0.62, SW - 1.24, SH - 1.24, None, ACCENT, 0.5)
    txt(s, 0, 1.7, SW, 0.4, "LUXURY WEDDING DESIGN & PRODUCTION", 14, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    rect(s, SW / 2 - 1.1, 2.35, 2.2, 0.025, ACCENT)
    txt(s, 0, 2.85, SW, 1.4, "SPARK PLANNERS", 60, TITLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    txt(s, 0, 4.35, SW, 0.6, "Designing unforgettable celebrations — plan · design · film · execute.",
        19, TEXT, font=HEAD, italic=True, align=PP_ALIGN.CENTER)
    rect(s, SW / 2 - 1.1, 5.35, 2.2, 0.025, ACCENT)
    txt(s, 0, 5.65, SW, 0.4, "A CLIENT PRESENTATION", 13, MUTED, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 0, 6.05, SW, 0.4, "by Sanjeeta Khaira", 15, TITLE, font=HEAD, italic=True, align=PP_ALIGN.CENTER)

    # 2 Who we are
    s = slide()
    kicker(s, "WHO WE ARE")
    title(s, "A full-service luxury wedding studio")
    txt(s, 0.92, 2.2, 6.4, 3.6,
        "Spark Planners is a boutique wedding design and production house led by "
        "Sanjeeta Khaira. We bring architecture, décor, film, and flawless on-ground "
        "execution under one roof — so every celebration is designed with precision "
        "and delivered with calm.\n\n"
        "From the first concept to the final frame, one team owns your story.",
        17, TEXT, spacing=1.15)
    for i, (h, d) in enumerate([("Design-led", "Every element drawn to scale before it is built."),
                                ("Cinematic", "In-house films capture the day, start to finish."),
                                ("Single team", "One accountable production team, end to end.")]):
        y = 2.2 + i * 1.15
        rect(s, 7.7, y, 4.7, 1.0, CARD, ACCENT, 0.75)
        txt(s, 7.95, y + 0.14, 4.2, 0.4, h, 16, TITLE, bold=True, font=HEAD)
        txt(s, 7.95, y + 0.55, 4.2, 0.4, d, 12.5, MUTED)
    footer(s)

    # 3 What we do
    s = slide(PANEL)
    kicker(s, "WHAT WE DO")
    title(s, "Six studios, one celebration")
    services = [
        ("Design Studio", "Venue, mandap, floral & lighting in 2D, 3D and VR."),
        ("Wedding Films", "Pre-wedding, highlights, drone and full films."),
        ("Production", "Timeline, crews, call sheets and on-ground delivery."),
        ("Guest & Hospitality", "RSVP, seating, VIP and welcome logistics."),
        ("Vendor Coordination", "Curated vendors across décor, lighting and film."),
        ("Vedi & Muhurat", "Vastu-aware vedi placement and auspicious windows."),
    ]
    cw, ch, gx = 3.78, 1.5, 0.32
    for i, (h, d) in enumerate(services):
        r, c = divmod(i, 3)
        x = 0.92 + c * (cw + gx); y = 2.15 + r * (ch + 0.34)
        rect(s, x, y, cw, ch, CARD, ACCENT, 0.75)
        rect(s, x, y, 0.08, ch, ACCENT)
        txt(s, x + 0.28, y + 0.18, cw - 0.5, 0.4, h, 15.5, TITLE, bold=True, font=HEAD)
        txt(s, x + 0.28, y + 0.62, cw - 0.5, 0.8, d, 12.5, MUTED)
    footer(s)

    # 4 Our difference
    s = slide()
    kicker(s, "OUR DIFFERENCE")
    title(s, "Powered by a Wedding Operating System")
    txt(s, 0.92, 2.0, 11.5, 0.7,
        "We don't sketch on napkins. Every Spark wedding is planned on a connected design platform — "
        "so what you approve is exactly what we build.", 16, TEXT, spacing=1.1)
    for i, (n, l) in enumerate([("156", "design surfaces"), ("1,441", "knowledge nodes"),
                                ("4,539", "linked relations"), ("17", "board pages")]):
        x = 0.92 + i * 3.0
        rect(s, x, 3.1, 2.7, 1.7, CARD, ACCENT, 0.75)
        txt(s, x, 3.35, 2.7, 0.9, n, 46, TITLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
        txt(s, x, 4.3, 2.7, 0.4, l, 12.5, MUTED, align=PP_ALIGN.CENTER)
    for i, t in enumerate(["AI design co-pilot for ideas & options", "3D / CAD studios for true-to-scale design",
                           "Board Composer for client-ready presentations", "Vedi & Muhurat intelligence built in"]):
        txt(s, 0.95, 5.2 + i * 0.42, 11.5, 0.4, "•  " + t, 14, TEXT)
    footer(s)

    # 5 How we work
    s = slide(PANEL)
    kicker(s, "HOW WE WORK")
    title(s, "From first call to final frame")
    steps = [("01", "Discover", "Vision, guest count, dates, budget."),
             ("02", "Design", "Venue, mandap, décor, lighting to scale."),
             ("03", "Muhurat & Vedi", "Vastu-aware placement & auspicious windows."),
             ("04", "Produce", "Crews, timeline, call sheets, logistics."),
             ("05", "Film", "Pre-wedding through full wedding film."),
             ("06", "Celebrate", "We run the day; you live it.")]
    for i, (n, h, d) in enumerate(steps):
        r, c = divmod(i, 3)
        x = 0.92 + c * 4.05; y = 2.2 + r * 1.9
        rect(s, x, y, 3.8, 1.6, CARD, ACCENT, 0.75)
        txt(s, x + 0.25, y + 0.16, 1.5, 0.6, n, 30, ACCENT, bold=True, font=HEAD)
        txt(s, x + 1.35, y + 0.22, 2.3, 0.4, h, 16, TITLE, bold=True, font=HEAD)
        txt(s, x + 0.25, y + 0.78, 3.3, 0.7, d, 12.5, MUTED)
    footer(s)

    # 5b Destination weddings
    s = slide()
    kicker(s, "DESTINATION WEDDINGS")
    title(s, "Weddings anywhere you dream")
    dests = ["spark/udaipur.png", "spark/jaipur.png", "spark/goa.png", "spark/himachal.png", "spark/thailand.png"]
    tw, gap = 2.156, 0.18
    for i, d in enumerate(dests):
        x = 0.92 + i * (tw + gap)
        pic_cover(s, G(d), x, 2.25, tw, 4.0)
    txt(s, 0.92, 6.4, 11.5, 0.4,
        "Palace, beach, mountain, or international — we scout the venue, manage guest travel and stay, and execute end to end.",
        13, MUTED, italic=True)
    footer(s)

    # 6 Proposal divider — dreamy full-bleed
    s = slide()
    _dream = prefer(G("spark/aisle.png"), G("udaipur-palace-night-hero.png"), P("extracted-boards/vedic-lotus-all-pages-board.png"))
    if _dream:
        pic_cover(s, _dream, 0, 0, SW, SH); overlay(s, 0, 0, SW, SH, 58)
    txt(s, 0, 1.75, SW, 0.4, "THE PROPOSAL", 13, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    rect(s, SW / 2 - 1.2, 2.4, 2.4, 0.025, ACCENT)
    txt(s, 0, 2.8, SW, 1.3, "A Dream in Ivory & Gold", 48, TITLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    txt(s, 0, 4.3, SW, 0.5, "Royal Heritage Wedding  ·  Udaipur Palace  ·  December 12, 2026", 17, TEXT, italic=True, font=HEAD, align=PP_ALIGN.CENTER)
    txt(s, 0, 5.0, SW, 0.5, "Where the lake holds the sky, and every light is a wish.", 14, MUTED, italic=True, align=PP_ALIGN.CENTER)

    # 7 Concept
    s = slide()
    kicker(s, "CONCEPT & VISION")
    title(s, "An evening that feels like a dream")
    txt(s, 0.92, 2.1, 6.1, 3.6,
        "Picture arriving by candlelit boat as the lake mirrors a sky of gold. A canopy of "
        "suspended lotus and orchids drifts above the saptapadi, and the palace glows like a "
        "lantern at dusk.\n\n"
        "We pair timeless Vedic ritual with quiet innovation — every detail drawn to scale, "
        "previewed in 3D, and lit to shift with your muhurat.",
        16, TEXT, spacing=1.18)
    txt(s, 0.95, 5.35, 6.1, 0.4, "Palette: Ivory · Gold · Lotus pink     Mood: candle-lit · ethereal · royal", 13, TITLE, italic=True)
    _concept_top = prefer(G("spark/invitation.png"), P("extracted-boards/vedic-sacred-design-options.png"))
    pic_cover(s, _concept_top, 7.45, 1.9, 4.95, 1.92)
    pic_cover(s, prefer(G("spark/carts.png"), _concept_top), 7.45, 3.98, 4.95, 1.92)
    footer(s)

    # 7b Signature experiences — innovative & dreamy moments
    s = slide(PANEL)
    kicker(s, "SIGNATURE EXPERIENCES")
    title(s, "Moments made to be remembered")
    exps = [
        ("Candlelit lake baraat", "Arrive by boat across the lake at golden hour."),
        ("Floating lotus canopy", "Suspended lotus & orchids drifting above the saptapadi."),
        ("Celestial lighting", "Uplighting that shifts to your sacred muhurat window."),
        ("Mirror saptapadi", "A reflecting pool ringed with floating diyas."),
        ("Scent design", "Jasmine entry · sandalwood mandap · rose dining."),
        ("Starlit finale", "A drone-lit night sky — subject to airspace clearance."),
    ]
    cw, ch, gx, gy = 3.78, 1.5, 0.32, 0.34
    for i, (h, d) in enumerate(exps):
        r, c = divmod(i, 3)
        x = 0.92 + c * (cw + gx); y = 2.15 + r * (ch + gy)
        rect(s, x, y, cw, ch, CARD, ACCENT, 0.75)
        rect(s, x, y, cw, 0.06, ACCENT)
        txt(s, x + 0.26, y + 0.22, cw - 0.5, 0.5, h, 15.5, TITLE, bold=True, font=HEAD)
        txt(s, x + 0.26, y + 0.72, cw - 0.5, 0.7, d, 12.5, MUTED, spacing=1.1)
    txt(s, 0.92, 6.5, 11.5, 0.4, "Designed for you — sensory, cinematic, and unmistakably yours.", 13, TEXT, italic=True)
    footer(s)

    # 8 Venue & layout
    s = slide(PANEL)
    kicker(s, "VENUE & LAYOUT")
    title(s, "Udaipur Palace — set for 620 guests")
    pic_cover(s, prefer(G("udaipur-venue-aerial.png"),
                        P("tlps-wedding-os/homepage/destinations/udaipur.png")), 7.55, 1.95, 4.85, 3.9)
    for i, (n, l) in enumerate([("96 × 64 m", "venue envelope"), ("620", "seated guests"),
                                ("3 / 6", "entries / exits"), ("4", "parking blocks")]):
        r, c = divmod(i, 2)
        x = 0.92 + c * 3.2; y = 2.2 + r * 1.55
        rect(s, x, y, 2.95, 1.35, CARD, ACCENT, 0.75)
        txt(s, x, y + 0.2, 2.95, 0.6, n, 28, TITLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
        txt(s, x, y + 0.92, 2.95, 0.3, l, 12, MUTED, align=PP_ALIGN.CENTER)
    txt(s, 0.95, 5.5, 6.3, 0.8, "Zones: entry court · mandap · stage · guest seating · dining · bride & groom rooms.", 13.5, TEXT)
    footer(s)

    # 9 Mandap
    s = slide()
    kicker(s, "MANDAP DESIGN")
    title(s, "Royal Heritage Mandap")
    _mandap = prefer(G("royal-heritage-mandap.png"), P("extracted-boards/vedic-sacred-traditional-options.png"))
    pic_cover(s, _mandap, 6.6, 1.95, 5.8, 4.35, focus=0.66)
    _mcap = ("Royal Heritage Mandap — design direction." if _mandap and "generated" in _mandap
             else "Sacred mandap design language — shortlisted directions.")
    txt(s, 6.6, 6.34, 5.8, 0.3, _mcap, 10.5, MUTED, italic=True)
    for i, (n, l) in enumerate([("12 × 7.5 m", "footprint"), ("8", "carved pillars"),
                                ("6.4 m", "canopy height"), ("5.2 m", "saptapadi dia")]):
        r, c = divmod(i, 2)
        x = 0.92 + c * 2.78; y = 2.15 + r * 1.5
        rect(s, x, y, 2.55, 1.3, CARD, ACCENT, 0.75)
        txt(s, x, y + 0.22, 2.55, 0.5, n, 23, TITLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
        txt(s, x, y + 0.82, 2.55, 0.3, l, 12, MUTED, align=PP_ALIGN.CENTER)
    txt(s, 0.92, 5.3, 5.5, 1.0, "Materials\nBrushed brass · lotus silk · ivory stone · neem greenery.", 13, TEXT, spacing=1.2)
    footer(s)

    # 10 Vedi & Muhurat
    s = slide(PANEL)
    kicker(s, "VEDI & MUHURAT")
    title(s, "Auspicious by design")
    for i, (h, v) in enumerate([("Vedi facing", "East"), ("Agni zone", "South-East"), ("Phera", "Clockwise")]):
        x = 0.92 + i * 3.9
        rect(s, x, 2.2, 3.6, 1.2, CARD, ACCENT, 0.75)
        txt(s, x + 0.25, 2.34, 3.1, 0.35, h, 13, MUTED, bold=True)
        txt(s, x + 0.25, 2.72, 3.1, 0.5, v, 22, TITLE, bold=True, font=HEAD)
    txt(s, 0.92, 3.75, 11.5, 0.4, "Recommended vivah muhurat windows", 15, ACCENT, bold=True, font=HEAD)
    muh = [("Uttara Phalguni – Panchami", "21:00–22:15", "97%"),
           ("Rohini – Tritiya", "19:10–20:30", "96%"),
           ("Uttara Ashadha – Ekadashi", "20:05–21:15", "93%")]
    for i, (name, win, score) in enumerate(muh):
        y = 4.25 + i * 0.72
        rect(s, 0.92, y, 11.5, 0.62, CARD, ACCENT, 0.5)
        txt(s, 1.15, y + 0.13, 6.5, 0.4, name, 15, TEXT, bold=True)
        txt(s, 7.8, y + 0.15, 2.2, 0.35, win, 13, MUTED)
        txt(s, 10.5, y + 0.12, 1.7, 0.4, score + " auspicious", 13, TITLE, bold=True)
    txt(s, 0.95, 6.55, 11.5, 0.35, "Indicative panchang windows — to be confirmed with the family priest.", 11, MUTED, italic=True)
    footer(s)

    # 11 Lighting & Floral
    s = slide()
    kicker(s, "LIGHTING & FLORAL")
    title(s, "Lit and dressed for the night")
    _lit = prefer(G("lighting-night-scene.png"))
    _flo = prefer(G("floral-lotus-marigold.png"))
    if _lit and _flo:
        pic_cover(s, _lit, 0.92, 2.05, 5.55, 2.55)
        pic_cover(s, _flo, 6.85, 2.05, 5.55, 2.55)
        txt(s, 0.92, 4.7, 5.55, 0.35, "LIGHTING", 14, ACCENT, bold=True)
        txt(s, 0.92, 5.1, 5.55, 1.6,
            "Central chandelier, warm stage wash and portal LED ribs.\n"
            "Warm 2700K key with cool accents; aisle path lights for the night.",
            13.5, TEXT, spacing=1.25)
        txt(s, 6.85, 4.7, 5.55, 0.35, "FLORAL", 14, ACCENT, bold=True)
        txt(s, 6.85, 5.1, 5.55, 1.6,
            "Lotus pink · marigold saffron · jasmine white · neem foliage.\n"
            "420 m of garland across canopy, aisle and entry toran.",
            13.5, TEXT, spacing=1.25)
    else:
        rect(s, 0.92, 2.2, 5.55, 4.05, CARD, ACCENT, 0.75)
        rect(s, 0.92, 2.2, 0.08, 4.05, ACCENT)
        rect(s, 6.85, 2.2, 5.55, 4.05, CARD, ACCENT, 0.75)
        rect(s, 6.85, 2.2, 0.08, 4.05, ACCENT)
        txt(s, 1.28, 2.5, 5.0, 0.4, "LIGHTING", 15, ACCENT, bold=True)
        txt(s, 1.28, 3.05, 5.0, 3.0,
            "Central chandelier over the mandap\n"
            "Warm stage wash for the ceremony\n"
            "Portal LED ribs framing the aisle\n"
            "Warm 2700K key with cool accents\n"
            "Aisle path lights for the night",
            14.5, TEXT, spacing=1.4)
        txt(s, 7.21, 2.5, 5.0, 0.4, "FLORAL", 15, ACCENT, bold=True)
        for i, (lab, col) in enumerate([("Lotus", LOTUS), ("Marigold", MARIGOLD),
                                        ("Jasmine", JASMINE), ("Neem", NEEM)]):
            cx = 7.21 + i * 1.27
            rect(s, cx, 3.05, 0.36, 0.36, col, ACCENT, 0.5, shape=MSO_SHAPE.OVAL)
            txt(s, cx - 0.2, 3.46, 0.76, 0.3, lab, 9.5, MUTED, align=PP_ALIGN.CENTER)
        txt(s, 7.21, 4.1, 5.0, 2.0,
            "Lotus-and-marigold canopy with jasmine accents and neem foliage framing the saptapadi.\n\n"
            "420 m of garland across canopy, aisle and entry toran.",
            14.5, TEXT, spacing=1.35)
    footer(s)

    # 12 Films
    s = slide(PANEL)
    kicker(s, "WEDDING FILMS")
    title(s, "Your story, on film")
    pic_cover(s, P("tlps-wedding-os/homepage/film-studio/full-wedding-films.png"), 7.4, 1.95, 5.0, 3.9)
    for i, (h, d) in enumerate([("Pre-Wedding Film", "Cinematic on-location story"),
                                ("Highlights", "The day, condensed beautifully"),
                                ("Drone Film", "Aerial of palace & mandap"),
                                ("Full Wedding Film", "Feature-length documentary edit")]):
        y = 2.2 + i * 1.0
        rect(s, 0.92, y, 6.1, 0.85, CARD, ACCENT, 0.6)
        txt(s, 1.15, y + 0.13, 5.6, 0.35, h, 15, TITLE, bold=True, font=HEAD)
        txt(s, 1.15, y + 0.5, 5.6, 0.3, d, 12, MUTED)
    footer(s)

    # 13 Investment
    s = slide()
    kicker(s, "INVESTMENT")
    title(s, "Indicative budget")
    for i, (l, v) in enumerate([("Mandap fabrication", "₹8,50,000"), ("Floral & foliage", "₹6,20,000"),
                                ("Lighting & truss", "₹7,30,000")]):
        y = 2.2 + i * 0.8
        rect(s, 0.92, y, 7.4, 0.66, CARD, ACCENT, 0.5)
        txt(s, 1.15, y + 0.15, 5.0, 0.4, l, 16, TEXT)
        txt(s, 6.0, y + 0.13, 2.1, 0.4, v, 17, TITLE, bold=True, align=PP_ALIGN.RIGHT)
    rect(s, 0.92, 4.7, 7.4, 0.8, ACCENT)
    txt(s, 1.15, 4.86, 4.0, 0.5, "Estimated total", 18, ON_ACCENT, bold=True, font=HEAD)
    txt(s, 5.6, 4.84, 2.5, 0.5, "₹22,00,000", 24, ON_ACCENT, bold=True, font=HEAD, align=PP_ALIGN.RIGHT)
    rect(s, 8.7, 2.2, 3.7, 3.3, PANEL, ACCENT, 0.75)
    txt(s, 8.95, 2.45, 3.2, 0.4, "WHAT'S INCLUDED", 12, ACCENT, bold=True)
    txt(s, 8.95, 2.9, 3.25, 2.5, "Design & 3D previews\nMandap, stage & décor build\nLighting & floral\nProduction & crews\nClient board package", 13.5, TEXT, spacing=1.25)
    txt(s, 0.95, 5.75, 11.4, 0.4, "Indicative estimate for scope shown; final quote follows venue walk-through and vendor confirmations.", 11.5, MUTED, italic=True)
    footer(s)

    # 14 Timeline
    s = slide(PANEL)
    kicker(s, "PRODUCTION TIMELINE")
    title(s, "Built with precision, ahead of the day")
    tl = [("Dec 10", "Stage deck & LED wall · site marking"),
          ("Dec 10–11", "Mandap structure & canopy"),
          ("Dec 11", "Lighting rig & focus"),
          ("Dec 12 · AM", "Floral dressing & final touch"),
          ("Dec 12 · PM", "Wedding day — we run it, you live it")]
    rect(s, 1.4, 2.4, 0.03, 3.7, ACCENT)
    for i, (d, t) in enumerate(tl):
        y = 2.4 + i * 0.78
        rect(s, 1.27, y + 0.05, 0.28, 0.28, ACCENT, shape=MSO_SHAPE.OVAL)
        txt(s, 1.85, y, 2.3, 0.4, d, 14, TITLE, bold=True, font=HEAD)
        txt(s, 4.3, y, 8.0, 0.5, t, 14.5, TEXT)
    footer(s)

    # 15 Next steps
    s = slide()
    kicker(s, "NEXT STEPS")
    title(s, "Let's create your story")
    for i, (n, h, d) in enumerate([("01", "Discovery call", "Align on vision, dates and budget."),
                                   ("02", "Refined proposal", "Tailored design, films and quote."),
                                   ("03", "Design & build", "3D previews, approvals, production."),
                                   ("04", "Celebrate", "Flawless delivery, captured on film.")]):
        x = 0.92 + i * 3.0
        rect(s, x, 2.4, 2.75, 2.6, CARD, ACCENT, 0.75)
        txt(s, x, 2.65, 2.75, 0.7, n, 34, ACCENT, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
        txt(s, x + 0.2, 3.5, 2.35, 0.4, h, 15, TITLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
        txt(s, x + 0.2, 3.95, 2.35, 0.9, d, 12.5, MUTED, align=PP_ALIGN.CENTER)
    txt(s, 0.92, 5.5, 11.5, 0.5, "Book a discovery call to begin.", 18, TEXT, font=HEAD, italic=True, align=PP_ALIGN.CENTER)
    footer(s)

    # 16 Thank you
    s = slide()
    if _hero:
        pic_cover(s, _hero, 0, 0, SW, SH); overlay(s, 0, 0, SW, SH, 68)
    rect(s, 0.45, 0.45, SW - 0.9, SH - 0.9, None, ACCENT, 1.25)
    rect(s, 0.62, 0.62, SW - 1.24, SH - 1.24, None, ACCENT, 0.5)
    rect(s, SW / 2 - 1.1, 2.55, 2.2, 0.025, ACCENT)
    txt(s, 0, 3.0, SW, 1.0, "Thank you", 52, TITLE, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    txt(s, 0, 4.2, SW, 0.5, "SPARK PLANNERS  ·  Sanjeeta Khaira", 18, TEXT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 0, 4.75, SW, 0.4, "Designing unforgettable celebrations — plan · design · film · execute.",
        13.5, MUTED, italic=True, align=PP_ALIGN.CENTER)
    rect(s, SW / 2 - 1.1, 5.5, 2.2, 0.025, ACCENT)

    prs.save(out_path)


def main():
    global prs, BLANK
    keys = sys.argv[1:] or list(THEMES)
    for key in keys:
        if key not in THEMES:
            print(f"unknown theme '{key}'. options: {', '.join(THEMES)}"); continue
        t = THEMES[key]
        apply_theme(t)
        prs = Presentation()
        prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
        BLANK = prs.slide_layouts[6]
        out = BASE + t["suffix"] + ".pptx"
        build_deck(out)
        print(f"saved {os.path.basename(out)}  ({t['label']}, {len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
